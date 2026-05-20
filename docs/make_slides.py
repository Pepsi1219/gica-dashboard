"""
Generate CSA Manager user guide as a PPTX.
Run: python make_slides.py
Output: csa-manager-guide.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette ────────────────────────────────────────────────────────────────
DEEP   = RGBColor(0x0B, 0x4E, 0x63)   # primary teal
BLUE   = RGBColor(0x3B, 0x8D, 0xC3)   # mid blue
INK    = RGBColor(0x0E, 0x28, 0x33)   # body text
MUTED  = RGBColor(0x5A, 0x70, 0x80)   # captions
LIGHT  = RGBColor(0xF4, 0xFA, 0xFB)   # very light bg
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xC9, 0xD8, 0xDE)
OK     = RGBColor(0x10, 0xB9, 0x81)
WARN   = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BLUE_S = RGBColor(0x3B, 0x82, 0xF6)   # training blue

FONT_TH = "Tahoma"   # safe Thai-capable font on Windows

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]

# ── helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Emu(9525)
    if not shadow:
        s.shadow.inherit = False
    return s

def add_rounded(slide, x, y, w, h, fill, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.10
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Emu(9525)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT_TH):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def bg(slide, color):
    add_rect(slide, 0, 0, SW, SH, color)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), DEEP)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11), Inches(0.55),
             title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.58), Inches(11), Inches(0.35),
                 subtitle, size=12, color=RGBColor(0xCA, 0xDC, 0xFC))
    # bottom-right pill
    pill_w, pill_h = Inches(1.6), Inches(0.4)
    pill_x = SW - pill_w - Inches(0.5)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   pill_x, Inches(0.25), pill_w, pill_h)
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    pill.line.fill.background()
    pill.shadow.inherit = False
    add_text(slide, pill_x, Inches(0.27), pill_w, Inches(0.36),
             "CSA Manager", size=11, bold=True, color=DEEP,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

def footer(slide, page):
    add_text(slide, Inches(0.6), SH - Inches(0.45), Inches(8), Inches(0.3),
             "New Operator Monitoring  ·  คู่มือผู้ใช้ CSA Manager",
             size=10, color=MUTED)
    add_text(slide, SW - Inches(1.2), SH - Inches(0.45), Inches(0.8), Inches(0.3),
             f"{page}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, DEEP)
# big decorative circle accent (top-right)
deco = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          SW - Inches(4), -Inches(2),
                          Inches(7), Inches(7))
deco.fill.solid(); deco.fill.fore_color.rgb = BLUE
deco.line.fill.background(); deco.shadow.inherit = False

# small pill label
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.8), Inches(1.4), Inches(2.4), Inches(0.5))
pill.adjustments[0] = 0.5
pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
pill.line.fill.background(); pill.shadow.inherit = False
add_text(s, Inches(0.8), Inches(1.42), Inches(2.4), Inches(0.46),
         "📘  คู่มือการใช้งาน", size=13, bold=True, color=DEEP,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.8), Inches(2.2), Inches(8.5), Inches(1.2),
         "CSA Manager", size=60, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.4), Inches(9), Inches(0.7),
         "โหมดผู้บริหาร · ดูข้อมูลอย่างเดียว",
         size=22, color=RGBColor(0xCA, 0xDC, 0xFC))

# Divider
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(0.8), Inches(4.5), Inches(0.6), Inches(0.06))
div.fill.solid(); div.fill.fore_color.rgb = WARN
div.line.fill.background(); div.shadow.inherit = False

add_text(s, Inches(0.8), Inches(4.7), Inches(11), Inches(0.5),
         "เข้าระบบได้ทันที · ไม่ต้องใช้รหัส · ดูข้อมูลทั้งหมด",
         size=17, color=WHITE)

add_text(s, Inches(0.8), Inches(6.7), Inches(8), Inches(0.4),
         "New Operator Monitoring System", size=11, color=RGBColor(0xCA, 0xDC, 0xFC))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — HOW TO LOGIN
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, LIGHT)
header_bar(s, "วิธีเข้าสู่ระบบ", "3 ขั้นตอน เข้าได้ทันทีโดยไม่ต้องใช้รหัส")

# 3 step cards
step_w = Inches(3.8); step_h = Inches(4.8)
step_y = Inches(1.5)
gap = Inches(0.3)
start_x = (SW - (step_w * 3 + gap * 2)) / 2

steps = [
    ("1", "เปิดเว็บไซต์ของระบบ",
     "ใส่ที่อยู่เว็บที่ทีมไอทีส่งให้ในเบราว์เซอร์\n(เช่น Chrome หรือ Edge)"),
    ("2", "คลิกปุ่ม “สำหรับ CSA Manager”",
     "ที่หน้าเข้าสู่ระบบ จะเห็น 2 ปุ่ม\nให้เลือกปุ่มล่าง"),
    ("3", "เข้าระบบสำเร็จ",
     "ระบบจะพาเข้าสู่หน้าแรกทันที\nไม่ต้องกรอก username/password"),
]

for i, (num, title, desc) in enumerate(steps):
    x = start_x + i * (step_w + gap)
    # card
    add_rounded(s, x, step_y, step_w, step_h, WHITE, BORDER)
    # number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                x + Inches(0.4), step_y + Inches(0.4),
                                Inches(1.0), Inches(1.0))
    circle.fill.solid(); circle.fill.fore_color.rgb = DEEP
    circle.line.fill.background(); circle.shadow.inherit = False
    add_text(s, x + Inches(0.4), step_y + Inches(0.4),
             Inches(1.0), Inches(1.0),
             num, size=36, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # title
    add_text(s, x + Inches(0.4), step_y + Inches(1.8),
             step_w - Inches(0.8), Inches(1.0),
             title, size=18, bold=True, color=DEEP)
    # description
    add_text(s, x + Inches(0.4), step_y + Inches(2.9),
             step_w - Inches(0.8), Inches(1.8),
             desc, size=14, color=INK)

# Hint box at bottom
hint_y = step_y + step_h + Inches(0.3)
add_rounded(s, Inches(1.0), hint_y, SW - Inches(2.0), Inches(0.55),
            RGBColor(0xFF, 0xF8, 0xD8), WARN)
add_text(s, Inches(1.3), hint_y + Inches(0.05),
         SW - Inches(2.6), Inches(0.5),
         "💡  เคล็ดลับ: ปุ่ม “สำหรับ CSA Manager” จะอยู่ใต้ปุ่ม Sign in with Microsoft เสมอ",
         size=12, color=INK, valign=MSO_ANCHOR.MIDDLE)

footer(s, 2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOME PAGE: SELECT DEPARTMENT
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, LIGHT)
header_bar(s, "หน้าแรก ① · เลือกหน่วยงาน",
           "ส่วนบนสุดของหน้าแรก แสดงหน่วยงานทั้งหมดให้เลือกเข้าดูข้อมูล")

# Left side: explanation
left_x = Inches(0.6)
add_text(s, left_x, Inches(1.3), Inches(5.5), Inches(0.5),
         "คุณจะเห็นอะไรบ้าง?", size=20, bold=True, color=DEEP)

bullets = [
    ("🏢", "กล่องของแต่ละหน่วยงาน",
     "G1, G2, G3, G4, TRM, EA"),
    ("👆", "คลิกกล่องไหน",
     "ระบบจะพาไปหน้า Dashboard ของหน่วยงานนั้น"),
    ("🔄", "เปลี่ยนหน่วยงานได้",
     "กดปุ่ม “← กลับหน้าแรก” ที่มุมซ้ายบนของหน้า Dashboard"),
]

y = Inches(1.9)
for icon, ttl, txt in bullets:
    add_text(s, left_x, y, Inches(0.5), Inches(0.5),
             icon, size=22, valign=MSO_ANCHOR.TOP)
    add_text(s, left_x + Inches(0.6), y, Inches(5), Inches(0.4),
             ttl, size=15, bold=True, color=INK)
    add_text(s, left_x + Inches(0.6), y + Inches(0.4),
             Inches(5), Inches(0.7),
             txt, size=13, color=MUTED)
    y += Inches(1.15)

# Right side: department grid mockup
grid_x = Inches(6.8)
grid_y = Inches(1.3)
add_rounded(s, grid_x, grid_y, Inches(6), Inches(5), WHITE, BORDER)
add_text(s, grid_x + Inches(0.3), grid_y + Inches(0.2),
         Inches(5.4), Inches(0.4),
         "เลือกหน่วยงาน", size=14, bold=True, color=BLUE)

names = ["G1", "G2", "G3", "G4", "TRM", "EA"]
cell_w = Inches(1.7); cell_h = Inches(1.2)
gap_x = Inches(0.15); gap_y = Inches(0.15)
sx = grid_x + Inches(0.3); sy = grid_y + Inches(0.8)
for i, n in enumerate(names):
    col = i % 3; row = i // 3
    cx = sx + col * (cell_w + gap_x)
    cy = sy + row * (cell_h + gap_y)
    btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              cx, cy, cell_w, cell_h)
    btn.adjustments[0] = 0.15
    btn.fill.solid(); btn.fill.fore_color.rgb = DEEP
    btn.line.fill.background(); btn.shadow.inherit = False
    add_text(s, cx, cy, cell_w, cell_h,
             n, size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

footer(s, 3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOME PAGE: OVERALL DASHBOARD
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, LIGHT)
header_bar(s, "หน้าแรก ② · ภาพรวมทั้งหมด",
           "ส่วนล่างของหน้าแรก สรุปข้อมูลของทุกหน่วยงานในที่เดียว")

# Top section: explanation
sections = [
    ("📊", "ตารางสรุป",
     "ดูจำนวนพนักงานทุกหน่วยงาน:\nทั้งหมด / ฝึกเสร็จ / กำลังฝึก /\nลาออก / โอนย้าย"),
    ("🍩", "Data Analytics",
     "กราฟวงกลม 4 ตัว + การ์ดสถิติ\n+ กราฟแนวโน้มรายเดือน"),
    ("🔘", "ฟิลเตอร์ปี",
     "เลือกปีเพื่อดูข้อมูลเฉพาะปี\nหรือ “ทุกปี” เพื่อดูภาพรวม"),
    ("💾", "ปุ่ม PNG / PDF",
     "บันทึกหน้านี้เป็นรูปภาพ\nหรือไฟล์ PDF ไว้ดูทีหลัง"),
]

# 4 cards in a row
card_w = Inches(3.0); card_h = Inches(2.5)
gap = Inches(0.2)
total_w = card_w * 4 + gap * 3
start_x = (SW - total_w) / 2
y = Inches(1.4)

for i, (icon, ttl, desc) in enumerate(sections):
    x = start_x + i * (card_w + gap)
    add_rounded(s, x, y, card_w, card_h, WHITE, BORDER)
    # icon circle
    icon_c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + Inches(0.3), y + Inches(0.3),
                                 Inches(0.7), Inches(0.7))
    icon_c.fill.solid(); icon_c.fill.fore_color.rgb = BLUE
    icon_c.line.fill.background(); icon_c.shadow.inherit = False
    add_text(s, x + Inches(0.3), y + Inches(0.3),
             Inches(0.7), Inches(0.7),
             icon, size=18, color=WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), y + Inches(1.1),
             card_w - Inches(0.6), Inches(0.4),
             ttl, size=16, bold=True, color=DEEP)
    add_text(s, x + Inches(0.3), y + Inches(1.55),
             card_w - Inches(0.6), Inches(1.0),
             desc, size=12, color=MUTED)

# Bottom highlight box
hl_y = Inches(4.3)
add_rounded(s, Inches(1.2), hl_y, SW - Inches(2.4), Inches(2.3),
            DEEP, None)

add_text(s, Inches(1.5), hl_y + Inches(0.25),
         SW - Inches(3), Inches(0.5),
         "เปลี่ยนฟิลเตอร์ → ทุกอย่างอัพเดททันที",
         size=20, bold=True, color=WHITE)
add_text(s, Inches(1.5), hl_y + Inches(0.9),
         SW - Inches(3), Inches(0.5),
         "เมื่อเลือกปีในช่องด้านขวาบน ตารางสรุปและกราฟทุกตัวจะอัพเดทพร้อมกัน",
         size=14, color=RGBColor(0xCA, 0xDC, 0xFC))
add_text(s, Inches(1.5), hl_y + Inches(1.5),
         SW - Inches(3), Inches(0.6),
         "💡 ภาพที่ export ออกมาจะตรงกับฟิลเตอร์ที่ตั้งไว้ทุกครั้ง",
         size=14, color=WARN)

footer(s, 4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DEPARTMENT DASHBOARD
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, LIGHT)
header_bar(s, "หน้า Dashboard ของหน่วยงาน",
           "เปิดดูได้เมื่อคลิกกล่องหน่วยงานที่หน้าแรก")

# 5 stat cards top
stat_w = Inches(2.35); stat_h = Inches(1.5)
gap = Inches(0.15)
total_w = stat_w * 5 + gap * 4
sx = (SW - total_w) / 2
sy = Inches(1.4)

stats = [
    ("ทั้งหมด", "0", INK, None),
    ("ฝึกเสร็จ", "0", OK, "ตรงเวลา · เกินกำหนด"),
    ("กำลังฝึก", "0", BLUE_S, "พื้นฐาน · ขั้นตอนงาน"),
    ("การลาออก", "0", DANGER, "ก่อน/ระหว่าง พื้นฐาน · ขั้นตอน"),
    ("การโอนย้าย", "0", PURPLE, "ก่อน/ระหว่าง พื้นฐาน · ขั้นตอน"),
]

for i, (lbl, val, color, sub) in enumerate(stats):
    x = sx + i * (stat_w + gap)
    add_rounded(s, x, sy, stat_w, stat_h, WHITE, BORDER)
    add_text(s, x, sy + Inches(0.15), stat_w, Inches(0.6),
             val, size=32, bold=True, color=color,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, x, sy + Inches(0.78), stat_w, Inches(0.3),
             lbl, size=12, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, x + Inches(0.15), sy + Inches(1.08),
                 stat_w - Inches(0.3), Inches(0.4),
                 sub, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Section: Below — Employee data card mockup + explanation
# Left: explanation
left_y = sy + stat_h + Inches(0.4)
add_text(s, Inches(0.6), left_y, Inches(6), Inches(0.5),
         "ส่วนต่างๆ ในหน้านี้", size=18, bold=True, color=DEEP)

items = [
    ("🔢", "การ์ดสรุปด้านบน 5 ใบ",
     "ดูจำนวนพนักงานในแต่ละสถานะของหน่วยงานนี้"),
    ("🔍", "ค้นหารหัสพนักงาน",
     "พิมพ์รหัสในช่องค้นหา จะเปิดหน้าต่างรายละเอียด"),
    ("🎚️", "ฟิลเตอร์ ปี / สถานะ",
     "เลือกปีหรือสถานะที่สนใจ ตารางจะเฉพาะรายการที่ตรง"),
    ("📋", "ตารางพนักงาน",
     "เลื่อนซ้าย-ขวาเพื่อดูทุกคอลัมน์ของแต่ละคน"),
]

iy = left_y + Inches(0.6)
for icon, ttl, desc in items:
    add_text(s, Inches(0.6), iy, Inches(0.5), Inches(0.4),
             icon, size=18)
    add_text(s, Inches(1.2), iy, Inches(5.5), Inches(0.35),
             ttl, size=13, bold=True, color=INK)
    add_text(s, Inches(1.2), iy + Inches(0.32), Inches(5.5), Inches(0.4),
             desc, size=11, color=MUTED)
    iy += Inches(0.75)

# Right: Mockup of search bar + status filter
mx = Inches(7.2); my = left_y
add_rounded(s, mx, my, Inches(5.5), Inches(3.5), WHITE, BORDER)

add_text(s, mx + Inches(0.3), my + Inches(0.2),
         Inches(5), Inches(0.4),
         "ข้อมูลพนักงาน", size=14, bold=True, color=DEEP)

# search input
add_rounded(s, mx + Inches(0.3), my + Inches(0.7),
            Inches(2.5), Inches(0.4),
            RGBColor(0xF4, 0xFA, 0xFB), BORDER)
add_text(s, mx + Inches(0.4), my + Inches(0.72),
         Inches(2.4), Inches(0.4),
         "🔍  พิมพ์รหัสพนักงาน...", size=10, color=MUTED, valign=MSO_ANCHOR.MIDDLE)
# year filter
add_rounded(s, mx + Inches(2.95), my + Inches(0.7),
            Inches(1.1), Inches(0.4),
            WHITE, BORDER)
add_text(s, mx + Inches(2.95), my + Inches(0.72),
         Inches(1.1), Inches(0.4),
         "ทุกปี ▾", size=10, color=INK,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
# status filter
add_rounded(s, mx + Inches(4.15), my + Inches(0.7),
            Inches(1.05), Inches(0.4),
            WHITE, BORDER)
add_text(s, mx + Inches(4.15), my + Inches(0.72),
         Inches(1.05), Inches(0.4),
         "สถานะ ▾", size=10, color=INK,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# table header bar
add_rect(s, mx + Inches(0.3), my + Inches(1.35),
         Inches(4.9), Inches(0.3), DEEP)
add_text(s, mx + Inches(0.4), my + Inches(1.37),
         Inches(4.7), Inches(0.3),
         "รหัส   ชื่อ       เกรด   เริ่มฝึก   ครบกำหนด   สถานะ",
         size=9, bold=True, color=WHITE)

# rows
row_data = [
    ("E001", "พนักงาน A",  "B", "ฝึกเสร็จ",         OK),
    ("E002", "พนักงาน B",  "C", "กำลังฝึก",          BLUE_S),
    ("E003", "พนักงาน C",  "D", "การลาออก",          DANGER),
    ("E004", "พนักงาน D",  "B", "การโอนย้าย",        PURPLE),
]
for i, (eid, nm, gr, st, c) in enumerate(row_data):
    ry = my + Inches(1.7) + i * Inches(0.4)
    bg_c = WHITE if i % 2 == 0 else RGBColor(0xFA, 0xFC, 0xFD)
    add_rect(s, mx + Inches(0.3), ry, Inches(4.9), Inches(0.35), bg_c)
    add_text(s, mx + Inches(0.4), ry + Inches(0.02),
             Inches(3), Inches(0.35),
             f"{eid}   {nm}   {gr}",
             size=10, color=INK, valign=MSO_ANCHOR.MIDDLE)
    # status badge
    add_rounded(s, mx + Inches(3.6), ry + Inches(0.06),
                Inches(1.5), Inches(0.24),
                RGBColor(0xFF, 0xFF, 0xFF), c)
    add_text(s, mx + Inches(3.6), ry + Inches(0.06),
             Inches(1.5), Inches(0.24),
             st, size=8.5, bold=True, color=c,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

footer(s, 5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PERMISSIONS (What you CAN and CANNOT do)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, LIGHT)
header_bar(s, "สิทธิ์การใช้งาน",
           "CSA Manager ดูข้อมูลได้อย่างเดียว · แก้ไขไม่ได้")

# Two columns
col_w = Inches(5.8); col_h = Inches(5)
col_y = Inches(1.4)
left_x = Inches(0.7)
right_x = SW - col_w - Inches(0.7)

# Left: CAN (green)
add_rounded(s, left_x, col_y, col_w, col_h, WHITE, OK)
# header strip
add_rect(s, left_x, col_y, col_w, Inches(0.7), OK)
add_text(s, left_x, col_y, col_w, Inches(0.7),
         "✓  ทำได้",
         size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

can_items = [
    "ดูตาราง “ภาพรวมทั้งหมด” ของทุกหน่วยงาน",
    "ดูสถิติและกราฟใน Data Analytics",
    "เปลี่ยนฟิลเตอร์ปีและสถานะตามใจ",
    "เปิด Dashboard ของแต่ละหน่วยงาน",
    "ค้นหารหัสพนักงาน → ดูรายละเอียด",
    "บันทึกหน้าเป็น PNG หรือ PDF",
    "เปลี่ยนภาษา (TH / EN / LO / VI)",
    "เปลี่ยนธีม สว่าง / มืด",
]
y = col_y + Inches(0.95)
for txt in can_items:
    add_text(s, left_x + Inches(0.3), y, Inches(0.4), Inches(0.4),
             "✓", size=15, bold=True, color=OK)
    add_text(s, left_x + Inches(0.75), y, col_w - Inches(1), Inches(0.4),
             txt, size=13, color=INK, valign=MSO_ANCHOR.MIDDLE)
    y += Inches(0.45)

# Right: CANNOT (red)
add_rounded(s, right_x, col_y, col_w, col_h, WHITE, DANGER)
add_rect(s, right_x, col_y, col_w, Inches(0.7), DANGER)
add_text(s, right_x, col_y, col_w, Inches(0.7),
         "✗  ทำไม่ได้",
         size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

cant_items = [
    "ลงทะเบียนพนักงานใหม่ (ปุ่ม “ลงทะเบียน” ถูกซ่อน)",
    "แก้ไขข้อมูลในตาราง (ปุ่ม ✏ ถูกซ่อน)",
    "บันทึก/อัพเดทข้อมูลในหน้าต่างค้นหา",
    "เปิดหน้า Admin (เฟือง ⚙ ถูกซ่อน)",
    "เพิ่ม/ลบวันหยุดพิเศษ",
]
y = col_y + Inches(0.95)
for txt in cant_items:
    add_text(s, right_x + Inches(0.3), y, Inches(0.4), Inches(0.4),
             "✗", size=15, bold=True, color=DANGER)
    add_text(s, right_x + Inches(0.75), y, col_w - Inches(1), Inches(0.45),
             txt, size=13, color=INK, valign=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)

# bottom note
add_text(s, right_x + Inches(0.3), y + Inches(0.2),
         col_w - Inches(0.6), Inches(0.4),
         "ระบบจะซ่อนปุ่มเหล่านี้ให้อัตโนมัติเมื่อเข้าผ่านโหมด CSA Manager",
         size=11, color=MUTED)

footer(s, 6)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TIPS & END
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, DEEP)

# Big rounded white card center
card_w = Inches(11); card_h = Inches(5.8)
card_x = (SW - card_w) / 2; card_y = (SH - card_h) / 2 - Inches(0.2)
add_rounded(s, card_x, card_y, card_w, card_h, WHITE)

# Title inside
add_text(s, card_x, card_y + Inches(0.5),
         card_w, Inches(0.7),
         "เคล็ดลับการใช้งาน", size=28, bold=True, color=DEEP,
         align=PP_ALIGN.CENTER)
add_text(s, card_x, card_y + Inches(1.2),
         card_w, Inches(0.4),
         "เพื่อให้ใช้งานได้สะดวกขึ้น", size=13, color=MUTED,
         align=PP_ALIGN.CENTER)

# 4 tip cards (2x2)
tips = [
    ("🌐", "เปลี่ยนภาษา",
     "คลิกธงด้านขวาบน เพื่อสลับ\nไทย / English / ลาว / Tiếng Việt"),
    ("🌙", "ธีมมืด / สว่าง",
     "คลิกไอคอนพระจันทร์/อาทิตย์\nบนแถบด้านขวาบน"),
    ("💾", "บันทึก PNG / PDF",
     "ในหน้าภาพรวม กดปุ่ม PNG/PDF\nไฟล์จะมีปีที่เลือกไว้แสดงด้วย"),
    ("📅", "ตัวกรองช่วงเวลา",
     "เปลี่ยน “ทุกปี” เป็นปีที่ต้องการ\nเพื่อดูข้อมูลของช่วงเวลานั้น"),
]
tw = Inches(4.8); th = Inches(1.7)
gx = Inches(0.3); gy = Inches(0.25)
sx = card_x + (card_w - (tw * 2 + gx)) / 2
sy = card_y + Inches(2.0)
for i, (icon, ttl, desc) in enumerate(tips):
    col = i % 2; row = i // 2
    x = sx + col * (tw + gx); y = sy + row * (th + gy)
    add_rounded(s, x, y, tw, th, LIGHT, None)
    # icon big
    add_text(s, x + Inches(0.2), y + Inches(0.2),
             Inches(0.8), Inches(0.8),
             icon, size=28, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.1), y + Inches(0.15),
             tw - Inches(1.3), Inches(0.45),
             ttl, size=15, bold=True, color=DEEP)
    add_text(s, x + Inches(1.1), y + Inches(0.65),
             tw - Inches(1.3), Inches(1.0),
             desc, size=11, color=MUTED)

# Bottom thank-you bar
add_text(s, 0, SH - Inches(0.7),
         SW, Inches(0.4),
         "สอบถามเพิ่มเติม · ติดต่อทีมไอที", size=14,
         color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.CENTER)

# ── Save ──
output = "csa-manager-guide.pptx"
prs.save(output)
print(f"Saved: {output}")
